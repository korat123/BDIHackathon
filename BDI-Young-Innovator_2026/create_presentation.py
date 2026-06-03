"""
Patient Rescue Radar — BDI Young Innovator 2026
Minimal presentation generator (python-pptx)
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0F, 0x17, 0x2A)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE   = RGBColor(0x25, 0x63, 0xEB)
C_BLUE2  = RGBColor(0x1E, 0x40, 0xAF)
C_ORANGE = RGBColor(0xF9, 0x73, 0x16)
C_RED    = RGBColor(0xEF, 0x44, 0x44)
C_GRAY   = RGBColor(0x64, 0x74, 0x8B)
C_LGRAY  = RGBColor(0xE2, 0xE8, 0xF0)
C_OFFWH  = RGBColor(0xF8, 0xFA, 0xFC)
C_DGRAY  = RGBColor(0x1E, 0x29, 0x3B)
C_ORNG_L = RGBColor(0xFF, 0xED, 0xD5)

# ── Slide geometry ───────────────────────────────────────────────
SW = Cm(33.87)
SH = Cm(19.05)
MX = Cm(2.0)
CW = SW - 2 * MX   # ~29.87 cm
FONT = "Segoe UI"


# ── Primitives ───────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, text="", size=14, bold=False,
        color=C_DARK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    r.font.italic = italic
    return tb


def mbox(slide, x, y, w, h, lines):
    """Multi-paragraph text box. Each line is a dict: text,size,bold,color,align,italic."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 14))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = ln.get("color", C_DARK)
        r.font.name = FONT
        r.font.italic = ln.get("italic", False)
    return tb


def rec(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.75)):
    sp = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO RECTANGLE
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = line
        sp.line.width = lw
    else:
        sp.line.fill.background()
    return sp


def hline(slide, x, y, w, color=C_LGRAY, h=Cm(0.05)):
    return rec(slide, x, y, w, h, fill=color)


def lbl(slide, text):
    box(slide, MX, Cm(0.85), CW, Cm(0.65),
        text=text, size=9, bold=True, color=C_BLUE)


def hed(slide, text, y=Cm(2.0), size=36, color=C_DARK, w=None, x=None):
    return box(slide, x or MX, y, w or CW, Cm(3.5),
               text=text, size=size, bold=True, color=color)


def cap(slide, text, color=C_GRAY):
    box(slide, MX, SH - Cm(1.5), CW, Cm(1.0),
        text=text, size=10, color=color, italic=True)


# ── Slide 1: Cover ───────────────────────────────────────────────

def s01_cover(prs, blank):
    sl = prs.slides.add_slide(blank)
    bg(sl, C_DARK)

    # Right-edge blue bar
    rec(sl, SW - Cm(0.45), Cm(0), Cm(0.45), SH, fill=C_BLUE)

    # Decorative accent lines top-right
    for i, (w_cm, col) in enumerate([(18, C_BLUE), (12, C_BLUE2), (7, C_DGRAY)]):
        rec(sl, SW - Cm(w_cm + 0.45), Cm(i * 0.28), Cm(w_cm), Cm(0.18), fill=col)

    # Event label
    box(sl, MX, Cm(1.6), Cm(20), Cm(0.7),
        text="BDI YOUNG INNOVATOR 2026", size=10, bold=True, color=C_BLUE)

    # Blue divider
    hline(sl, MX, Cm(2.6), Cm(13), color=C_BLUE, h=Cm(0.07))

    # Main title (two lines)
    mbox(sl, MX, Cm(3.6), Cm(24), Cm(7.0), [
        dict(text="Patient", size=76, bold=True, color=C_WHITE),
        dict(text="Rescue Radar", size=76, bold=True, color=C_WHITE),
    ])

    # Tagline
    box(sl, MX, Cm(10.2), Cm(26), Cm(1.4),
        text="AI ที่เปลี่ยน Reactive Care เป็น Proactive Care",
        size=18, color=C_GRAY)

    # Team footer
    hline(sl, MX, SH - Cm(2.6), Cm(11), color=C_BLUE2, h=Cm(0.06))
    box(sl, MX, SH - Cm(2.2), Cm(24), Cm(1.1),
        text="ทีม KMUTT  |  Health Track  |  BDI Hackathon 2026",
        size=12, color=C_GRAY)


# ── Slide 2: The Problem ─────────────────────────────────────────

def s02_problem(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "THE PROBLEM")
    hed(sl, "23 ล้านคน. ดูแลได้ดีแค่ 11.7%", y=Cm(2.2), size=38)
    hline(sl, MX, Cm(4.4), CW)

    BW = (CW - Cm(1.0)) / 3
    BH = Cm(8.5)
    BY = Cm(5.2)
    GAP = Cm(0.5)

    data = [
        ("23M",     C_BLUE,   "ผู้ป่วย NCD ในไทย",      "เบาหวาน + ความดัน"),
        ("88.3%",   C_ORANGE, "ไม่ได้รับการดูแลที่ดี",  "Inadequate NCD control"),
        ("20–37%",  C_RED,    "LTFU ทุกปี",             "หายไปจากระบบ"),
    ]
    for i, (num, nc, lbl_t, sub) in enumerate(data):
        x = MX + i * (BW + GAP)
        rec(sl, x, BY, BW, BH, fill=C_OFFWH)
        # Color top bar
        rec(sl, x, BY, BW, Cm(0.45), fill=nc)
        box(sl, x, BY + Cm(1.0), BW, Cm(3.2),
            text=num, size=60, bold=True, color=nc, align=PP_ALIGN.CENTER)
        box(sl, x, BY + Cm(4.3), BW, Cm(1.1),
            text=lbl_t, size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        box(sl, x, BY + Cm(5.5), BW, Cm(0.9),
            text=sub, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    cap(sl, "Source: Aekplakorn et al. 2025  |  WHO Thailand NCD Report")


# ── Slide 3: The Silent Cascade ──────────────────────────────────

def s03_cascade(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "THE SILENT CRISIS")
    hed(sl, "ระบบรู้ช้าเกินไป", y=Cm(2.0), size=36)
    box(sl, MX, Cm(4.1), CW, Cm(0.9),
        text="เส้นทางผู้ป่วย NCD จากวันวินิจฉัย ไปจนถึง ICU", size=14, color=C_GRAY)
    hline(sl, MX, Cm(5.2), CW)

    BASE_Y = Cm(10.5)
    hline(sl, MX, BASE_Y, CW, color=C_LGRAY, h=Cm(0.07))

    nodes = [
        (Cm(4.0),  "วินิจฉัย",        C_BLUE,   "Diagnosed"),
        (Cm(10.5), "รู้สึกดีขึ้น",     C_GRAY,   "Feels Better"),
        (Cm(17.0), "หยุดมาพบแพทย์",   C_ORANGE, "Stops Visiting"),
        (Cm(23.5), "ค่าพุ่งกลับ",      C_RED,    "BP/Glucose Spike"),
        (Cm(30.0), "STROKE / ICU",    C_RED,    "Emergency"),
    ]

    # Gap zone (orange zone between node 3 and 4)
    rec(sl, Cm(17.5), BASE_Y - Cm(1.8), Cm(5.5), Cm(3.6), fill=C_ORNG_L)
    box(sl, Cm(17.5), BASE_Y - Cm(1.6), Cm(5.5), Cm(0.9),
        text="6–14 เดือน", size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    box(sl, Cm(17.5), BASE_Y - Cm(0.6), Cm(5.5), Cm(0.7),
        text="ระบบมองไม่เห็น", size=11, color=C_ORANGE, align=PP_ALIGN.CENTER)
    box(sl, Cm(17.5), BASE_Y + Cm(1.0), Cm(5.5), Cm(0.7),
        text="invisible gap", size=10, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    NR = Cm(0.9)
    for nx, name_th, color, name_en in nodes:
        rec(sl, nx - NR/2, BASE_Y - NR/2, NR, NR, fill=color)
        box(sl, nx - Cm(2.2), BASE_Y - Cm(3.5), Cm(4.4), Cm(1.5),
            text=name_th, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        box(sl, nx - Cm(2.2), BASE_Y + Cm(1.2), Cm(4.4), Cm(0.8),
            text=name_en, size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

    cap(sl, "ค่า ICU ≈ 200,000–500,000 บาท  |  ป้องกันได้หากรู้ก่อน 6 เดือน  |  นี่คือปัญหาที่ AI แก้ได้")


# ── Slide 4: The Insight ─────────────────────────────────────────

def s04_insight(prs, blank):
    sl = prs.slides.add_slide(blank)
    bg(sl, C_DARK)
    box(sl, MX, Cm(0.85), Cm(20), Cm(0.65),
        text="OUR CORE INSIGHT", size=9, bold=True, color=C_BLUE)

    # Left column headlines
    LEFT_W = Cm(14.0)
    box(sl, MX, Cm(2.5), LEFT_W, Cm(2.2),
        text="ข้อมูลที่หายไป  ≠  Noise", size=36, bold=True, color=C_WHITE)
    box(sl, MX, Cm(5.0), LEFT_W, Cm(2.2),
        text="ข้อมูลที่หายไป  =  สัญญาณ", size=36, bold=True, color=C_ORANGE)

    hline(sl, MX, Cm(7.5), LEFT_W, color=C_DGRAY, h=Cm(0.06))

    box(sl, MX, Cm(8.1), LEFT_W, Cm(1.3),
        text="เมื่อผู้ป่วยไม่มา → ไม่มีค่าบันทึก\nนั่นคือข้อมูลที่สำคัญที่สุด",
        size=13, color=C_GRAY, wrap=True)

    box(sl, MX, SH - Cm(1.5), Cm(24), Cm(0.9),
        text="Informative Missingness (MNAR)  —  Sperrin et al. 2021, JMIR",
        size=10, color=RGBColor(0x3B, 0x52, 0x6B), italic=True)

    # EMR grid (right side)
    GX = Cm(18.0)
    GY0 = Cm(2.0)
    CW_G = Cm(1.65)
    CH_G = Cm(1.4)
    GAP  = Cm(0.22)

    COL_LABELS = ["SBP", "DBP", "HbA1c", "FPG", "Creat.", "Meds", "Visit", "BMI"]
    PATTERN = [
        [True,  False, True,  False, False, True,  False, False],
        [False, True,  False, False, True,  False, True,  False],
        [True,  False, False, True,  False, False, False, True ],
        [False, False, True,  False, True,  False, True,  False],
    ]
    ROW_LABELS = ["P−1", "P0", "P1", "P2"]

    # Column headers
    for c, cl in enumerate(COL_LABELS):
        box(sl, GX + c * (CW_G + GAP), GY0 - Cm(0.75), CW_G, Cm(0.7),
            text=cl, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    for r, rl in enumerate(ROW_LABELS):
        box(sl, GX - Cm(2.1), GY0 + r * (CH_G + GAP), Cm(1.9), CH_G,
            text=rl, size=10, bold=True, color=C_GRAY, align=PP_ALIGN.RIGHT)
        for c in range(8):
            filled = PATTERN[r][c]
            fx = GX + c * (CW_G + GAP)
            fy = GY0 + r * (CH_G + GAP)
            fc = C_BLUE if filled else C_DGRAY
            lc = C_BLUE2 if filled else RGBColor(0x2D, 0x3F, 0x5A)
            rec(sl, fx, fy, CW_G, CH_G, fill=fc, line=lc, lw=Pt(0.5))
            if not filled:
                box(sl, fx, fy + Cm(0.15), CW_G, CH_G - Cm(0.15),
                    text="—", size=18, color=RGBColor(0x2D, 0x3F, 0x5A),
                    align=PP_ALIGN.CENTER)


# ── Slide 5: Pipeline ────────────────────────────────────────────

def s05_pipeline(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "THE SOLUTION")
    hed(sl, "Patient Rescue Radar", y=Cm(2.2), size=34)
    hline(sl, MX, Cm(4.3), CW)

    BW = (CW - Cm(1.5)) / 4
    BH = Cm(8.0)
    BY = Cm(5.2)
    GAP = Cm(0.5)

    steps = [
        (C_DARK,   "1", "220K+ EMR",          "ข้อมูล EMR ผู้ป่วย\nเบาหวาน + ความดัน\n(DM 70K + HT 150K)"),
        (C_BLUE,   "2", "Feature Engineering","12 signals จาก\n4,000+ คอลัมน์\n(sparsity + visit gap)"),
        (C_BLUE2,  "3", "XGBoost Model",      "เรียนรู้ pattern\nการหายไปจากระบบ\n(sparse-native)"),
        (C_ORANGE, "4", "Risk Score + SHAP",  "0–100 + 3 เหตุผล\nที่อ่านได้ทันที\n(PDPA-ready)"),
    ]

    for i, (fill, num, title, sub) in enumerate(steps):
        x = MX + i * (BW + GAP)
        rec(sl, x, BY, BW, BH, fill=fill)

        # Step badge
        rec(sl, x + Cm(0.5), BY + Cm(0.5), Cm(0.8), Cm(0.8), fill=C_WHITE)
        box(sl, x + Cm(0.5), BY + Cm(0.45), Cm(0.8), Cm(0.85),
            text=num, size=12, bold=True, color=fill, align=PP_ALIGN.CENTER)

        box(sl, x + Cm(0.4), BY + Cm(1.7), BW - Cm(0.8), Cm(1.3),
            text=title, size=15, bold=True, color=C_WHITE, wrap=True)
        box(sl, x + Cm(0.4), BY + Cm(3.2), BW - Cm(0.8), Cm(4.0),
            text=sub, size=12, color=RGBColor(0xC7, 0xD7, 0xF4) if fill != C_ORANGE else RGBColor(0xFF, 0xF0, 0xE0),
            wrap=True)

        # Connector
        if i < 3:
            cx = x + BW
            cy = BY + BH / 2
            hline(sl, cx, cy - Cm(0.06), GAP, color=C_LGRAY, h=Cm(0.12))

    cap(sl, "Prototype validated on 200-patient sample  |  Scales to 220K+ patients on full dataset")


# ── Slide 6: Patient Card ────────────────────────────────────────

def s06_output(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "WHAT THE OUTPUT LOOKS LIKE")
    hed(sl, "Actionable, Explainable, Ready to Act", y=Cm(2.1), size=28)
    hline(sl, MX, Cm(4.0), CW)

    CX = MX + Cm(0.5)
    CY = Cm(4.7)
    CRDW = Cm(19.5)
    CRDH = Cm(12.8)

    # Card body
    rec(sl, CX, CY, CRDW, CRDH, fill=C_DARK)

    # Header bar
    rec(sl, CX, CY, CRDW, Cm(2.1), fill=C_DGRAY)
    # Red dot
    rec(sl, CX + Cm(0.6), CY + Cm(0.6), Cm(0.9), Cm(0.9), fill=C_RED)
    box(sl, CX + Cm(1.9), CY + Cm(0.45), Cm(9.0), Cm(0.8),
        text="สมชาย อ.  |  67 ปี  |  เบาหวาน Type 2",
        size=13, bold=True, color=C_WHITE)

    # Risk badge
    rec(sl, CX + Cm(13.0), CY + Cm(0.35), Cm(6.0), Cm(1.4), fill=C_RED)
    box(sl, CX + Cm(13.0), CY + Cm(0.4), Cm(6.0), Cm(1.2),
        text="Risk Score: 87/100  HIGH",
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Divider
    hline(sl, CX + Cm(0.6), CY + Cm(2.3), CRDW - Cm(1.2), color=C_DGRAY)

    # Warnings
    warnings = [
        "ไม่มีการนัดหมาย 14 เดือน",
        "ไม่มีบันทึก HbA1c ใน 2 ปีที่ผ่านมา",
        "BP ล่าสุด: 158 mmHg (ยังไม่ถูกควบคุม)",
    ]
    for j, w in enumerate(warnings):
        wy = CY + Cm(2.8) + j * Cm(2.3)
        rec(sl, CX + Cm(0.6), wy + Cm(0.25), Cm(0.55), Cm(0.55), fill=C_ORANGE)
        box(sl, CX + Cm(1.5), wy + Cm(0.05), CRDW - Cm(2.1), Cm(1.0),
            text=w, size=14, bold=True, color=C_WHITE)

    # Recommendation
    hline(sl, CX + Cm(0.6), CY + Cm(9.7), CRDW - Cm(1.2), color=C_DGRAY)
    rec(sl, CX, CY + Cm(9.9), CRDW, Cm(2.9), fill=C_BLUE)
    box(sl, CX + Cm(0.7), CY + Cm(10.1), CRDW - Cm(1.2), Cm(1.5),
        text="→  แนะนำ: Case Manager โทรหาภายใน 7 วัน",
        size=16, bold=True, color=C_WHITE)

    # Right legend
    LX = CX + CRDW + Cm(1.2)
    LW = MX + CW - LX

    box(sl, LX, Cm(5.0), LW, Cm(0.9),
        text="Risk Threshold", size=13, bold=True, color=C_DARK)
    rec(sl, LX, Cm(6.1), Cm(0.9), Cm(0.9), fill=C_RED)
    box(sl, LX + Cm(1.2), Cm(6.1), LW - Cm(1.4), Cm(1.5),
        text="Score ≥ 70\nActive Intervention\n(Case Manager + อสม.)", size=11, color=C_DARK)
    rec(sl, LX, Cm(8.3), Cm(0.9), Cm(0.9), fill=C_BLUE)
    box(sl, LX + Cm(1.2), Cm(8.3), LW - Cm(1.4), Cm(1.0),
        text="Score < 70  Routine Care", size=11, color=C_DARK)

    hline(sl, LX, Cm(9.9), LW, color=C_LGRAY)

    box(sl, LX, Cm(10.3), LW, Cm(3.5),
        text="SHAP-based:\nทุก decision มีเหตุผล\nที่อ่านได้ → PDPA-ready",
        size=11, color=C_GRAY, wrap=True)


# ── Slide 7: Impact ──────────────────────────────────────────────

def s07_impact(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "THE IMPACT")
    hline(sl, MX, Cm(1.9), CW)

    # Giant ROI
    box(sl, MX, Cm(2.4), CW, Cm(5.2),
        text="134×", size=120, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    box(sl, MX, Cm(7.4), CW, Cm(1.1),
        text="ผลตอบแทนจากการลงทุน (Conservative ROI Estimate)",
        size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    hline(sl, MX, Cm(9.0), CW)

    cols = [
        ("2M บาท",          C_DARK,   "เงินลงทุน",         "Development + Deployment"),
        ("270–540M บาท/ปี", C_BLUE,   "ประหยัดต่อปี",      "จากการลดค่าฟอกไต"),
        ("15–20%",          C_ORANGE, "ลด LTFU ปีแรก",     "Target Year 1"),
    ]
    CW3 = (CW - Cm(1.0)) / 3
    for i, (num, color, lbl_t, sub) in enumerate(cols):
        x = MX + i * (CW3 + Cm(0.5))
        box(sl, x, Cm(9.8), CW3, Cm(2.5),
            text=num, size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
        box(sl, x, Cm(12.4), CW3, Cm(1.0),
            text=lbl_t, size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        box(sl, x, Cm(13.5), CW3, Cm(0.9),
            text=sub, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    cap(sl, "Dialysis cost: 540K บาท/คน/ปี  |  220K patients eligible  |  1.6 trillion baht/year NCD burden")


# ── Slide 8: Infrastructure Ready ───────────────────────────────

def s08_infra(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "THAILAND IS READY")
    hed(sl, "โครงสร้างพื้นฐานมีอยู่แล้ว\nAI เพียงแค่ Activate", y=Cm(2.0), size=30)
    hline(sl, MX, Cm(5.6), CW)

    infra = [
        ("10,000+", C_BLUE,   "รพ.สต.",          "Health Centers\nทั่วประเทศไทย"),
        ("1.04M",   C_ORANGE, "อสม.",             "Volunteer Health Workers\nพร้อมลงพื้นที่"),
        ("2023–27", C_DARK,   "Smart NCD Policy", "นโยบายรัฐบาล\nรองรับ AI ใน NCD"),
    ]
    CW3 = (CW - Cm(1.0)) / 3
    for i, (num, color, lbl_t, sub) in enumerate(infra):
        x = MX + i * (CW3 + Cm(0.5))
        y = Cm(6.5)
        rec(sl, x, y, Cm(0.4), Cm(5.5), fill=color)
        box(sl, x + Cm(0.8), y, CW3 - Cm(1.0), Cm(2.5),
            text=num, size=46, bold=True, color=color)
        box(sl, x + Cm(0.8), y + Cm(2.6), CW3 - Cm(1.0), Cm(1.0),
            text=lbl_t, size=14, bold=True, color=C_DARK)
        box(sl, x + Cm(0.8), y + Cm(3.7), CW3 - Cm(1.0), Cm(2.0),
            text=sub, size=12, color=C_GRAY, wrap=True)

    cap(sl, "ไม่ต้องสร้าง infrastructure ใหม่ — ส่ง risk list ให้ถูกคนถูกที่ในเวลาที่ใช่")


# ── Slide 9: Technical + Ethics ──────────────────────────────────

def s09_tech(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "BUILT TO BE TRUSTED")
    hed(sl, "Technical Credibility & Ethics", y=Cm(2.2), size=30)
    hline(sl, MX, Cm(4.1), CW)

    COL_W = (CW - Cm(2.0)) / 2
    LX = MX
    RX = MX + COL_W + Cm(2.0)

    # Vertical divider
    rec(sl, MX + COL_W + Cm(0.9), Cm(4.5), Cm(0.06), Cm(12.0), fill=C_LGRAY)

    box(sl, LX, Cm(4.6), COL_W, Cm(0.9),
        text="Technical Stack", size=15, bold=True, color=C_BLUE)
    tech = [
        "XGBoost — sparse-native, fast, interpretable",
        "12 features engineered จาก 4,000+ EMR columns",
        "SHAP explainability → human-readable reasons",
        "On-premise only — ข้อมูลไม่ออกนอก รพ.",
        "AUROC benchmark: 85.9% (literature reference)",
        "Prototype 200 คน verified → scale 220K+ ได้",
    ]
    for j, item in enumerate(tech):
        y = Cm(5.8) + j * Cm(1.6)
        rec(sl, LX, y + Cm(0.3), Cm(0.25), Cm(0.8), fill=C_BLUE)
        box(sl, LX + Cm(0.6), y, COL_W - Cm(0.8), Cm(1.5), text=item, size=12, color=C_DARK)

    box(sl, RX, Cm(4.6), COL_W, Cm(0.9),
        text="PDPA & Fairness", size=15, bold=True, color=C_ORANGE)
    ethics = [
        "Output = risk score เท่านั้น — ไม่เปิดเผยข้อมูลดิบ",
        "SHAP = right to explanation (ตาม PDPA)",
        "Fairness audit แยกตาม age group & region",
        "ไม่ใช้ demographic features (gender/ethnicity)",
        "Audit trail — ทุก prediction มี timestamp + version",
    ]
    for j, item in enumerate(ethics):
        y = Cm(5.8) + j * Cm(1.6)
        rec(sl, RX, y + Cm(0.3), Cm(0.25), Cm(0.8), fill=C_ORANGE)
        box(sl, RX + Cm(0.6), y, COL_W - Cm(0.8), Cm(1.5), text=item, size=12, color=C_DARK)


# ── Slide 10: Team ───────────────────────────────────────────────

def s10_team(prs, blank):
    sl = prs.slides.add_slide(blank)
    lbl(sl, "OUR TEAM")
    hed(sl, "KMUTT Computer Engineering — Year 1", y=Cm(2.2), size=26)
    hline(sl, MX, Cm(4.2), CW)

    members = [
        ("กอ",      "Technical &\nAI Architect",     C_BLUE),
        ("แบงค์",   "Business &\nDomain Lead",       C_DARK),
        ("จีน",     "System Design\n& Workflow",     C_DARK),
        ("ข้าวฟาง", "Product &\nInnovation",         C_DARK),
        ("เฟิม",    "Data Science\n& Research",      C_DARK),
    ]
    CARDW = (CW - Cm(1.0)) / 5
    CARDH = Cm(8.5)
    CARDY = Cm(5.0)
    for i, (name, role, color) in enumerate(members):
        x = MX + i * (CARDW + Cm(0.25))
        rec(sl, x, CARDY, CARDW, CARDH, fill=C_OFFWH, line=C_LGRAY, lw=Pt(0.75))
        rec(sl, x, CARDY, CARDW, Cm(0.55), fill=color)
        box(sl, x, CARDY + Cm(1.0), CARDW, Cm(2.6),
            text=name, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        box(sl, x + Cm(0.2), CARDY + Cm(3.8), CARDW - Cm(0.4), Cm(3.5),
            text=role, size=12, color=C_GRAY, align=PP_ALIGN.CENTER, wrap=True)

    box(sl, MX, Cm(14.2), CW, Cm(1.3),
        text="จุฬาภรณราชวิทยาลัย alumni  |  Hackathon experience  |  Engineering & Health Data Science",
        size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 11: Closing ────────────────────────────────────────────

def s11_closing(prs, blank):
    sl = prs.slides.add_slide(blank)
    bg(sl, C_DARK)

    # Same accent bars as cover
    rec(sl, SW - Cm(0.45), Cm(0), Cm(0.45), SH, fill=C_BLUE)
    for i, (w_cm, col) in enumerate([(18, C_BLUE), (12, C_BLUE2), (7, C_DGRAY)]):
        rec(sl, SW - Cm(w_cm + 0.45), Cm(i * 0.28), Cm(w_cm), Cm(0.18), fill=col)

    box(sl, MX, Cm(1.6), Cm(20), Cm(0.65),
        text="BDI YOUNG INNOVATOR 2026", size=9, bold=True, color=C_BLUE)
    hline(sl, MX, Cm(2.6), Cm(13), color=C_BLUE, h=Cm(0.07))

    mbox(sl, MX, Cm(3.8), Cm(28), Cm(7.5), [
        dict(text="จากนี้ไป", size=48, bold=True, color=C_GRAY),
        dict(text="เราจะรู้ก่อน", size=64, bold=True, color=C_WHITE),
        dict(text="ที่ผู้ป่วยจะหายไป", size=64, bold=True, color=C_ORANGE),
    ])

    box(sl, MX, Cm(11.2), Cm(26), Cm(1.2),
        text="Patient Rescue Radar  —  AI-powered LTFU Prevention for NCD",
        size=16, color=C_GRAY)

    hline(sl, MX, SH - Cm(2.6), Cm(11), color=C_BLUE2, h=Cm(0.06))
    box(sl, MX, SH - Cm(2.2), Cm(20), Cm(1.0),
        text="ทีม KMUTT  |  BDI Hackathon 2026",
        size=12, color=RGBColor(0x3B, 0x52, 0x6B))


# ── Main ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]   # Blank layout

    s01_cover(prs, blank)
    s02_problem(prs, blank)
    s03_cascade(prs, blank)
    s04_insight(prs, blank)
    s05_pipeline(prs, blank)
    s06_output(prs, blank)
    s07_impact(prs, blank)
    s08_infra(prs, blank)
    s09_tech(prs, blank)
    s10_team(prs, blank)
    s11_closing(prs, blank)

    out = "BDI-Young-Innovator_2026/PatientRescueRadar_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
